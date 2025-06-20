import os
import tempfile
import warnings
from typing import List, Dict, Any, Optional, Tuple
import logging
import json
import base64
from io import BytesIO
from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv

# Document processing
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_community.document_loaders import UnstructuredPowerPointLoader, UnstructuredExcelLoader

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# PowerPoint image extraction
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. Install with: pip install python-pptx")

# Word document image extraction
try:
    from docx import Document as DocxDocument
    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml
    import zipfile
    DOCX_IMAGE_SUPPORT = True
except ImportError:
    DOCX_IMAGE_SUPPORT = False
    print("Warning: python-docx image support not available")

# Image processing and OCR
import torch
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
import fitz  # PyMuPDF for better PDF text extraction with coordinates
import cv2
import numpy as np

from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI

# Reranking
import cohere

# LLM - Azure OpenAI instead of Gemini
from langchain.prompts import ChatPromptTemplate

# Configure logging with file output
import os
from datetime import datetime

# Create logs directory if it doesn't exist
os.makedirs("logs", exist_ok=True)

# Create a timestamp for log files
timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
log_filename = f"logs/rag_pipeline_{timestamp}.log"

# Configure logging with both file and console output
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s | %(levelname)-8s | %(name)-15s | %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S',
    handlers=[
        logging.FileHandler(log_filename, encoding='utf-8'),
        logging.StreamHandler()  # Console output
    ]
)

logger = logging.getLogger(__name__)
logger.info("="*80)
logger.info("RAG PIPELINE LOGGING STARTED")
logger.info(f"Log file: {log_filename}")
logger.info("="*80)

# Suppress warnings
warnings.filterwarnings("ignore")

load_dotenv()

class ImageAnalyzer:
    """Analyze images/tables/infographics in PDF pages using Azure OpenAI Vision"""
    
    def __init__(self):
        # Initialize Azure OpenAI Vision LLM
        self.vision_llm = AzureChatOpenAI(
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT"),
            api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
            temperature=0
        )
    

    def analyze_image_with_llm(self, image: Image.Image):
        """
        Use Azure OpenAI Vision to analyze and summarize image content
        """
        try:
            # Convert image to base64 for the API
            import io
            import base64
            
            buffer = io.BytesIO()
            image.save(buffer, format='PNG')
            image_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            # Create the message with image
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": """Analyze this image and provide a detailed summary of its content. Include:
                                    1. Any text visible in the image
                                    2. Tables, charts, or data visualizations
                                    3. Key information, numbers, or insights
                                    4. Overall context and purpose

                                    Provide a comprehensive summary that captures all important information:"""
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_base64}"
                            }
                        }
                    ]
                }
            ]
            
            response = self.vision_llm.invoke(messages)
            return response.content if hasattr(response, 'content') else str(response)
            
        except Exception as e:
            logger.error(f"Error analyzing image with Azure OpenAI: {e}")
            return f"Image content detected but could not be analyzed: {str(e)}"

class DocumentProcessor:
    """Enhanced document processor with conditional text + image chunking"""
    
    def __init__(self, image_analyzer: ImageAnalyzer):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        self.image_analyzer = image_analyzer
        # Thresholds for determining text sufficiency
        self.min_text_length = 100  # Minimum characters for sufficient text
        self.min_meaningful_words = 10  # Minimum meaningful words
    
    def _is_text_sufficient(self, text: str):
        """
        Check if extracted text is sufficient for meaningful content
        
        Args:
            text: Extracted text to evaluate
            
        Returns:
            True if text is sufficient, False if image processing needed
        """
        if not text or not text.strip():
            return False
        
        # Remove extra whitespace and clean text
        cleaned_text = ' '.join(text.split())
        
        # Check minimum length
        if len(cleaned_text) < self.min_text_length:
            return False
        
        # Check for meaningful words (not just numbers, symbols, or very short words)
        words = cleaned_text.split()
        meaningful_words = [word for word in words if len(word) > 2 and word.isalpha()]
        
        if len(meaningful_words) < self.min_meaningful_words:
            return False
        
        logger.info(f"Text sufficient: {len(cleaned_text)} chars, {len(meaningful_words)} meaningful words")
        return True
    
    def _has_actual_images_in_pdf_page(self, page):
        """
        Check if a PDF page actually contains embedded images
        
        Args:
            page: PyMuPDF page object
            
        Returns:
            True if page contains actual images, False otherwise
        """
        try:
            # Get image list from the page
            image_list = page.get_images(full=True)
            
            if not image_list:
                return False
            
            # Check if images are substantial (not just tiny icons or decorative elements)
            substantial_images = 0
            for img_index, img in enumerate(image_list):
                try:
                    # Get image info
                    xref = img[0]
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    # Check image dimensions - filter out very small images (likely icons)
                    if pix.width > 50 and pix.height > 50:
                        substantial_images += 1
                    
                    pix = None  # Free memory
                    
                except Exception as e:
                    logger.debug(f"Error checking image {img_index}: {e}")
                    continue
            
            logger.info(f"Found {substantial_images} substantial images in page")
            return substantial_images > 0
            
        except Exception as e:
            logger.error(f"Error checking for images in PDF page: {e}")
            return False
    
    def _has_images_present(self, file_path: str, page=None, slide=None, docx_doc=None, xlsx_file=None):
        """
        Check if images are present in the content
        
        Args:
            file_path: Path to the file
            page: PDF page object (optional)
            slide: PowerPoint slide object (optional)
            docx_doc: Word document object (optional)
            xlsx_file: Excel file path (optional)
            
        Returns:
            True if images are detected, False otherwise
        """
        try:
            if page is not None:
                # For PDF: check if page actually has embedded images
                return self._has_actual_images_in_pdf_page(page)
            
            elif slide is not None:
                # For PowerPoint: check if slide has image shapes
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        return True
                return False
            
            elif docx_doc is not None:
                # For Word: check if images exist in the zip archive using advanced method
                images = self._extract_images_from_docx_advanced(file_path)
                return len(images) > 0
            
            elif xlsx_file is not None:
                # For Excel: check if workbook has embedded images
                try:
                    # Try to import openpyxl
                    try:
                        from openpyxl import load_workbook
                    except ImportError:
                        logger.debug("openpyxl not available for Excel image detection")
                        return False
                    
                    # Quick check for images in Excel file
                    workbook = load_workbook(xlsx_file)
                    has_images = False
                    
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        if hasattr(sheet, '_images') and sheet._images:
                            has_images = True
                            break
                    
                    workbook.close()
                    return has_images
                    
                except Exception as e:
                    logger.debug(f"Error checking Excel images: {e}")
                    return False
            
            return False
            
        except Exception as e:
            logger.error(f"Error checking for images: {e}")
            return False
    
    def process_documents(self, file_paths: List[str]):
        """
        Process documents with hybrid text + image chunking
        
        Returns:
            Tuple of (documents, visual_metadata)
        """
        all_documents = []
        visual_metadata = {}
        
        logger.info("="*60)
        logger.info("DOCUMENT PROCESSING STARTED")
        logger.info(f"Total documents to process: {len(file_paths)}")
        logger.info("="*60)
        
        for idx, file_path in enumerate(file_paths, 1):
            doc_name = os.path.basename(file_path)
            logger.info(f"[{idx}/{len(file_paths)}] Processing: {doc_name}")
            logger.info(f"   File path: {file_path}")
            logger.info(f"   File size: {os.path.getsize(file_path) / 1024:.1f} KB")
            
            start_time = datetime.now()
            
            try:
                if file_path.lower().endswith('.pdf'):
                    docs, metadata = self._process_pdf_conditional(file_path)
                    all_documents.extend(docs)
                    visual_metadata[file_path] = metadata
                else:
                    # For non-PDF files, use standard loaders
                    docs = self._load_non_pdf(file_path)
                    all_documents.extend(docs)
                    visual_metadata[file_path] = {"type": "non_pdf", "pages": []}
                    
                    # Log summary for non-PDF files
                    file_extension = os.path.splitext(file_path)[1].lower()
                    logger.info(f"   Document Summary for {doc_name}:")
                    logger.info(f"      File type: {file_extension.upper()[1:]} document")
                    logger.info(f"      Total chunks generated: {len(docs)}")
                    if file_extension in ['.pptx', '.docx', '.xlsx']:
                        logger.info(f"      Content processed: Mixed text and potential visual elements")
                    else:
                        logger.info(f"      Content processed: Text only")
                
                processing_time = (datetime.now() - start_time).total_seconds()
                logger.info(f"   SUCCESS: {doc_name} processed in {processing_time:.2f}s")
                logger.info(f"   Generated {len(docs)} chunks from {doc_name}")
                logger.info(f"   Total chunks so far: {len(all_documents)}")
                
            except Exception as e:
                processing_time = (datetime.now() - start_time).total_seconds()
                logger.error(f"   FAILED: {doc_name} failed in {processing_time:.2f}s")
                logger.error(f"   Error: {str(e)}")
                continue
            
            logger.info("-" * 40)
        
        logger.info("="*60)
        logger.info("DOCUMENT PROCESSING COMPLETED")
        logger.info(f"Successfully processed: {len([f for f in file_paths if f in visual_metadata])} documents")
        logger.info(f"Total chunks generated: {len(all_documents)}")
        logger.info(f"Documents with visual content: {len([m for m in visual_metadata.values() if m.get('type') == 'pdf'])}")
        logger.info("="*60)
        
        return all_documents, visual_metadata
    
    def _process_pdf_conditional(self, file_path: str):
        """Process PDF with conditional text + image chunking"""
        documents = []
        doc_name = os.path.basename(file_path)
        
        try:
            # First, extract text using PyMuPDF (fast)
            logger.info(f"   Opening PDF: {doc_name}")
            pdf_doc = fitz.open(file_path)
            logger.info(f"   PDF has {len(pdf_doc)} pages")
            
            pages_metadata = []
            page_images = None  # Only convert if needed
            
            for page_num in range(len(pdf_doc)):
                logger.info(f"   Processing {doc_name} - Page {page_num + 1}/{len(pdf_doc)}")
                page = pdf_doc[page_num]
                page_text = page.get_text()
                
                # Store basic page metadata
                page_metadata = {
                    "page_num": page_num,
                    "image": None,  # Will be populated only if needed
                    "text_blocks": None,
                    "full_text": page_text
                }
                
                page_chunks = []
                
                # Check text sufficiency and image presence
                text_sufficient = self._is_text_sufficient(page_text)
                has_images = self._has_images_present(file_path, page=page)
                
                # Track content types found on this page
                page_content_types = []
                
                # Always process text if sufficient
                if text_sufficient:
                    # Process text chunks
                    text_chunks = self.text_splitter.split_text(page_text)
                    
                    for i, chunk in enumerate(text_chunks):
                        doc = Document(
                            page_content=chunk,
                            metadata={
                                "source": file_path,
                                "page": page_num,
                                "chunk_id": f"{file_path}_page_{page_num}_text_{i}",
                                "chunk_index": i,
                                "content_type": "text"
                            }
                        )
                        page_chunks.append(doc)
                    
                    page_content_types.append(f"text ({len(text_chunks)} chunks)")
                
                # Always extract embedded images if present (regardless of text sufficiency)
                if has_images:
                    # Extract embedded images directly from PDF page
                    extracted_images = self._extract_images_from_pdf_page(page, page_num)
                    
                    if extracted_images:
                        # Analyze extracted embedded images
                        source_info = {
                            "source": file_path,
                            "page": page_num,
                            "chunk_index": len(page_chunks)
                        }
                        image_chunks = self._analyze_extracted_images(extracted_images, source_info)
                        page_chunks.extend(image_chunks)
                        page_content_types.append(f"images ({len(image_chunks)} summaries)")
                
                # Log what was found on this page
                if page_content_types:
                    content_summary = " and ".join(page_content_types)
                    logger.info(f"      Page {page_num + 1}: Found {content_summary}")
                else:
                    logger.info(f"      Page {page_num + 1}: No processable content found")
                
                # If no sufficient text and no embedded images found, fallback to page-level image analysis
                if not text_sufficient and has_images and not any(chunk.metadata.get("content_type") == "image_summary" for chunk in page_chunks):
                    logger.info(f"Page {page_num}: No embedded images found, using page-level image analysis")
                    
                    # Fallback: Convert entire page to image
                    if page_images is None:
                        logger.info("Converting PDF pages to images (fallback for page-level analysis)")
                        page_images = convert_from_path(file_path, dpi=200)
                        logger.info(f"Extracted {len(page_images)} page images")
                    
                    # Get page image and text blocks with coordinates
                    page_image = page_images[page_num] if page_num < len(page_images) else None
                    text_blocks = page.get_text("dict")
                    
                    # Update metadata
                    page_metadata["image"] = page_image
                    page_metadata["text_blocks"] = text_blocks
                    
                    if page_image:
                        # Analyze the page image with Vision LLM
                        logger.info(f"Analyzing page {page_num} with Vision LLM")
                        image_summary = self.image_analyzer.analyze_image_with_llm(page_image)
                        
                        if image_summary and image_summary.strip():
                            doc = Document(
                                page_content=f"[IMAGE SUMMARY] {image_summary}",
                                metadata={
                                    "source": file_path,
                                    "page": page_num,
                                    "chunk_id": f"{file_path}_page_{page_num}_image_0",
                                    "chunk_index": len(page_chunks),
                                    "content_type": "image_summary",
                                    "image_type": "page_image"
                                }
                            )
                            page_chunks.append(doc)
                            logger.info(f"Created page-level image summary for page {page_num}")
                
                # If no sufficient text, include any available text as fallback
                if not text_sufficient and page_text.strip():
                    text_chunks = self.text_splitter.split_text(page_text)
                    for i, chunk in enumerate(text_chunks):
                        doc = Document(
                            page_content=chunk,
                            metadata={
                                "source": file_path,
                                "page": page_num,
                                "chunk_id": f"{file_path}_page_{page_num}_text_{i}",
                                "chunk_index": len(page_chunks),
                                "content_type": "text"
                            }
                        )
                        page_chunks.append(doc)
                
                # Final fallback: if no content at all, use any available text
                if not page_chunks and page_text.strip():
                    logger.info(f"Page {page_num}: No sufficient content found, using available text as final fallback")
                    text_chunks = self.text_splitter.split_text(page_text)
                    for i, chunk in enumerate(text_chunks):
                        doc = Document(
                            page_content=chunk,
                            metadata={
                                "source": file_path,
                                "page": page_num,
                                "chunk_id": f"{file_path}_page_{page_num}_text_{i}",
                                "chunk_index": i,
                                "content_type": "text"
                            }
                        )
                        page_chunks.append(doc)
                
                pages_metadata.append(page_metadata)
                documents.extend(page_chunks)
            
            pdf_doc.close()
            
            # Generate document summary
            total_pages = len(pages_metadata)
            pages_with_text = 0
            pages_with_images = 0
            pages_with_both = 0
            total_text_chunks = 0
            total_image_chunks = 0
            
            for page_num in range(total_pages):
                page_docs = [doc for doc in documents if doc.metadata.get("page") == page_num]
                text_docs = [doc for doc in page_docs if doc.metadata.get("content_type") == "text"]
                image_docs = [doc for doc in page_docs if doc.metadata.get("content_type") == "image_summary"]
                
                has_text = len(text_docs) > 0
                has_images = len(image_docs) > 0
                
                if has_text and has_images:
                    pages_with_both += 1
                elif has_text:
                    pages_with_text += 1
                elif has_images:
                    pages_with_images += 1
                
                total_text_chunks += len(text_docs)
                total_image_chunks += len(image_docs)
            
            # Log document summary
            logger.info(f"   Document Summary for {doc_name}:")
            logger.info(f"      Total pages: {total_pages}")
            logger.info(f"      Pages with text only: {pages_with_text}")
            logger.info(f"      Pages with images only: {pages_with_images}")
            logger.info(f"      Pages with both text and images: {pages_with_both}")
            logger.info(f"      Total text chunks: {total_text_chunks}")
            logger.info(f"      Total image chunks: {total_image_chunks}")
            logger.info(f"      Total chunks generated: {len(documents)}")
            
            metadata = {
                "type": "pdf",
                "pages": pages_metadata
            }
            
            return documents, metadata
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            import traceback
            logger.error(f"Full traceback: {traceback.format_exc()}")
            raise
    
    def _load_non_pdf(self, file_path: str):
        """Load non-PDF documents with conditional processing"""
        if file_path.lower().endswith('.pptx'):
            return self._process_powerpoint_conditional(file_path)
        elif file_path.lower().endswith('.docx'):
            return self._process_docx_conditional(file_path)
        elif file_path.lower().endswith('.txt'):
            return self._process_document_conditional(file_path, TextLoader, "Text")
        elif file_path.lower().endswith('.xlsx'):
            return self._process_xlsx_conditional(file_path)
        elif file_path.lower().endswith('.csv'):
            # Convert CSV to XLSX first, then process
            xlsx_path = self._convert_csv_to_xlsx(file_path)
            if xlsx_path:
                return self._process_xlsx_conditional(xlsx_path)
            else:
                # Fallback: process CSV as text
                return self._process_document_conditional(file_path, TextLoader, "CSV")
        else:
            raise ValueError(f"Unsupported file format: {file_path}")
    
    def _process_document_conditional(self, file_path: str, loader_class, doc_type: str):
        """Process simple documents (txt, xlsx) with basic text extraction"""
        try:
            logger.info(f"Processing {doc_type} file: {file_path}")
            
            # Load document using the specified loader
            loader = loader_class(file_path)
            documents = loader.load()
            
            if not documents:
                logger.warning(f"No content found in {doc_type} file")
                return []
            
            # Combine all document content
            full_text = "\n\n".join([doc.page_content for doc in documents])
            
            # Split into chunks
            text_chunks = self.text_splitter.split_text(full_text)
            
            # Create document chunks
            chunks = []
            for i, chunk in enumerate(text_chunks):
                doc_chunk = Document(
                    page_content=chunk,
                    metadata={
                        "source": file_path,
                        "chunk_id": f"{file_path}_{doc_type.lower()}_{i}",
                        "chunk_index": i,
                        "page": i + 1,
                        "content_type": "text"
                    }
                )
                chunks.append(doc_chunk)
            
            logger.info(f"Successfully processed {doc_type} file: {len(chunks)} chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing {doc_type} file {file_path}: {e}")
            return []
        
    def _process_docx_conditional(self, file_path: str):
        """Process .docx files with conditional text + image extraction"""
        try:
            logger.info(f"Processing .docx file: {file_path}")
            
            if not DOCX_IMAGE_SUPPORT:
                logger.warning("python-docx image support not available, falling back to text-only extraction")
                return self._process_document_conditional(file_path, Docx2txtLoader, "Word")
            
            # Extract text using python-docx
            docx_doc = DocxDocument(file_path)
            
            # Extract all text content
            doc_text = ""
            for paragraph in docx_doc.paragraphs:
                doc_text += paragraph.text + "\n"
            
            # Also extract text from tables
            for table in docx_doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        doc_text += cell.text + " "
                doc_text += "\n"
            
            chunks = []
            
            # Check text sufficiency and image presence
            text_sufficient = self._is_text_sufficient(doc_text)
            has_images = self._has_images_present(file_path, docx_doc=docx_doc)
            
            # Always process text if sufficient
            if text_sufficient:
                logger.info(f".docx file: Text is sufficient, processing text chunks")
                
                # Process text chunks
                text_chunks = self.text_splitter.split_text(doc_text)
                
                for i, chunk in enumerate(text_chunks):
                    doc_chunk = Document(
                        page_content=chunk,
                        metadata={
                            "source": file_path,
                            "chunk_id": f"{file_path}_text_{i}",
                            "chunk_index": i,
                            "page": i + 1,
                            "content_type": "text"
                        }
                    )
                    chunks.append(doc_chunk)
            
            # Always extract embedded images if present (regardless of text sufficiency)
            if has_images:
                logger.info(f".docx file: Extracting embedded images")
                
                # Extract images using advanced method
                extracted_images = self._extract_images_from_docx_advanced(file_path)
                
                if extracted_images:
                    # Analyze extracted embedded images
                    source_info = {
                        "source": file_path,
                        "page": 1,
                        "chunk_index": len(chunks)
                    }
                    image_chunks = self._analyze_extracted_images(extracted_images, source_info)
                    chunks.extend(image_chunks)
                    logger.info(f"Added {len(image_chunks)} extracted image summaries for .docx file")
            
            # Fallback: if text insufficient and no images, still process available text
            elif not text_sufficient and doc_text.strip():
                logger.info(f".docx file: Text insufficient, using available text as fallback")
                
                text_chunks = self.text_splitter.split_text(doc_text)
                for i, chunk in enumerate(text_chunks):
                    doc_chunk = Document(
                        page_content=chunk,
                        metadata={
                            "source": file_path,
                            "chunk_id": f"{file_path}_text_{i}",
                            "chunk_index": len(chunks),
                            "page": len(chunks) + i + 1,
                            "content_type": "text"
                        }
                    )
                    chunks.append(doc_chunk)
            
            logger.info(f"Successfully processed .docx file: {len(chunks)} chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing .docx file {file_path}: {str(e)}")
            # Fallback to text-only processing
            logger.info("Falling back to text-only processing")
            return self._process_document_conditional(file_path, Docx2txtLoader, "Word")
    
    def _process_powerpoint_conditional(self, file_path: str):
        """Process PowerPoint files with conditional text + image extraction"""
        documents = []
        
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available, falling back to text-only extraction")
            loader = UnstructuredPowerPointLoader(file_path)
            return self._process_fallback_pptx(loader, file_path)
        
        try:
            presentation = Presentation(file_path)
            logger.info(f"Processing PowerPoint with {len(presentation.slides)} slides")
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                logger.info(f"Processing slide {slide_num}")
                
                # Extract text from slide
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                slide_chunks = []
                
                # Check text sufficiency and image presence
                text_sufficient = self._is_text_sufficient(slide_text)
                has_images = self._has_images_present(file_path, slide=slide)
                
                # Always process text if sufficient
                if text_sufficient:
                    logger.info(f"Slide {slide_num}: Text is sufficient, processing text chunks")
                    
                    # Process text chunks
                    text_chunks = self.text_splitter.split_text(slide_text)
                    for i, chunk in enumerate(text_chunks):
                        doc = Document(
                            page_content=chunk,
                            metadata={
                                "source": file_path,
                                "page": slide_num,
                                "chunk_id": f"{file_path}_slide_{slide_num}_text_{i}",
                                "chunk_index": len(documents),
                                "content_type": "text"
                            }
                        )
                        slide_chunks.append(doc)
                
                # Always extract embedded images if present (regardless of text sufficiency)
                if has_images:
                    logger.info(f"Slide {slide_num}: Extracting embedded images")
                    
                    # Extract embedded images directly from slide
                    extracted_images = self._extract_images_from_pptx_slide(slide, slide_num)
                    
                    if extracted_images:
                        # Analyze extracted embedded images
                        source_info = {
                            "source": file_path,
                            "page": slide_num,
                            "chunk_index": len(documents) + len(slide_chunks)
                        }
                        image_chunks = self._analyze_extracted_images(extracted_images, source_info)
                        slide_chunks.extend(image_chunks)
                        logger.info(f"Added {len(image_chunks)} extracted image summaries for slide {slide_num}")
                
                # Fallback: if text insufficient and no images, still process available text
                elif not text_sufficient and slide_text.strip():
                    logger.info(f"Slide {slide_num}: Text insufficient, using available text as fallback")
                    
                    text_chunks = self.text_splitter.split_text(slide_text)
                    for i, chunk in enumerate(text_chunks):
                        doc = Document(
                            page_content=chunk,
                            metadata={
                                "source": file_path,
                                "page": slide_num,
                                "chunk_id": f"{file_path}_slide_{slide_num}_text_{i}",
                                "chunk_index": len(documents) + len(slide_chunks),
                                "content_type": "text"
                            }
                        )
                        slide_chunks.append(doc)
                
                documents.extend(slide_chunks)
                logger.info(f"Slide {slide_num}: Created {len(slide_chunks)} chunks")
            
            logger.info(f"Successfully processed PowerPoint: {len(documents)} total chunks from {len(presentation.slides)} slides")
            return documents
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {str(e)}")
            # Fallback to text-only extraction
            loader = UnstructuredPowerPointLoader(file_path)
            return self._process_fallback_pptx(loader, file_path)

    def _process_xlsx_conditional(self, file_path: str):
        """Process Excel files with conditional text + image extraction"""
        try:
            logger.info(f"Processing .xlsx file: {file_path}")
            
            # Extract text using UnstructuredExcelLoader
            loader = UnstructuredExcelLoader(file_path)
            documents = loader.load()
            
            if not documents:
                logger.warning("No content found in Excel file")
                return []
            
            # Combine all document content
            doc_text = "\n\n".join([doc.page_content for doc in documents])
            
            chunks = []
            
            # Check text sufficiency and image presence
            text_sufficient = self._is_text_sufficient(doc_text)
            has_images = self._has_images_present(file_path, xlsx_file=file_path)
            
            # Always process text if sufficient
            if text_sufficient:
                logger.info(f".xlsx file: Text is sufficient, processing text chunks")
                
                # Process text chunks
                text_chunks = self.text_splitter.split_text(doc_text)
                
                for i, chunk in enumerate(text_chunks):
                    doc_chunk = Document(
                        page_content=chunk,
                        metadata={
                            "source": file_path,
                            "chunk_id": f"{file_path}_text_{i}",
                            "chunk_index": i,
                            "page": i + 1,
                            "content_type": "text"
                        }
                    )
                    chunks.append(doc_chunk)
            
            # Always extract embedded images if present (regardless of text sufficiency)
            if has_images:
                logger.info(f".xlsx file: Extracting embedded images")
                
                # Extract images using openpyxl method
                extracted_images = self._extract_images_from_xlsx(file_path)
                
                if extracted_images:
                    # Analyze extracted embedded images
                    source_info = {
                        "source": file_path,
                        "page": 1,
                        "chunk_index": len(chunks)
                    }
                    image_chunks = self._analyze_extracted_images(extracted_images, source_info)
                    chunks.extend(image_chunks)
                    logger.info(f"Added {len(image_chunks)} extracted image summaries for .xlsx file")
            
            # Fallback: if text insufficient and no images, still process available text
            elif not text_sufficient and doc_text.strip():
                logger.info(f".xlsx file: Text insufficient, using available text as fallback")
                
                text_chunks = self.text_splitter.split_text(doc_text)
                for i, chunk in enumerate(text_chunks):
                    doc_chunk = Document(
                        page_content=chunk,
                        metadata={
                            "source": file_path,
                            "chunk_id": f"{file_path}_text_{i}",
                            "chunk_index": len(chunks),
                            "page": len(chunks) + i + 1,
                            "content_type": "text"
                        }
                    )
                    chunks.append(doc_chunk)
            
            logger.info(f"Successfully processed .xlsx file: {len(chunks)} chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing .xlsx file {file_path}: {str(e)}")
            # Fallback to basic text-only processing
            logger.info("Falling back to basic Excel text processing")
            return self._process_document_conditional(file_path, UnstructuredExcelLoader, "Excel")

    def _extract_images_from_pdf_page(self, page, page_num):
        images = []
        
        try:
            # Get list of images on the page
            image_list = page.get_images(full=True)
            if image_list:
                logger.debug(f"Found {len(image_list)} embedded images on page {page_num + 1}")
            
            for img_index, img in enumerate(image_list):
                try:
                    # Get image reference
                    xref = img[0]
                    
                    # Extract image data
                    pix = fitz.Pixmap(page.parent, xref)
                    
                    # Skip very small images (likely icons or decorative elements)
                    if pix.width < 50 or pix.height < 50:
                        logger.debug(f"Skipping small image {img_index}: {pix.width}x{pix.height}")
                        pix = None
                        continue
                    
                    # Convert to PIL Image
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("ppm")
                        pil_image = Image.open(BytesIO(img_data))
                    else:  # CMYK
                        pix_rgb = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix_rgb.tobytes("ppm")
                        pil_image = Image.open(BytesIO(img_data))
                        pix_rgb = None
                    
                    images.append(pil_image)
                    logger.debug(f"Extracted image {img_index} from page {page_num + 1}: {pil_image.size}")
                    
                    # Clean up
                    pix = None
                    
                except Exception as e:
                    logger.error(f"Error extracting image {img_index} from page {page_num + 1}: {e}")
                    continue
                    
        except Exception as e:
            logger.error(f"Error getting images from page {page_num + 1}: {e}")
        
        return images
    
    def _extract_images_from_pptx_slide(self, slide, slide_num):
        images = []
        
        try:
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Extract image data
                        image_stream = shape.image.blob
                        pil_image = Image.open(BytesIO(image_stream))
                        
                        # Skip very small images
                        if pil_image.width < 50 or pil_image.height < 50:
                            logger.debug(f"Skipping small image {shape_index}: {pil_image.size}")
                            continue
                        
                        images.append(pil_image)
                        logger.debug(f"Extracted image {shape_index} from slide {slide_num}: {pil_image.size}")
                        
                    except Exception as e:
                        logger.error(f"Error extracting image {shape_index} from slide {slide_num}: {e}")
                        continue
                        
        except Exception as e:
            logger.error(f"Error getting images from slide {slide_num}: {e}")
        
        return images
    
    def _extract_images_from_docx_advanced(self, file_path):
        images = []
        
        try:
            # Method 1: Extract from ZIP archive (most reliable)
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                # List all files in the archive
                file_list = docx_zip.namelist()
                
                # Find image files in various locations
                image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']
                image_files = []
                
                # Check common image locations
                for file_name in file_list:
                    if any(file_name.lower().endswith(ext) for ext in image_extensions):
                        # Images can be in different folders
                        if (file_name.startswith('word/media/') or 
                            file_name.startswith('word/embeddings/') or
                            file_name.startswith('media/')):
                            image_files.append(file_name)
                
                logger.info(f"Found {len(image_files)} image files in .docx archive")
                
                for img_index, image_file in enumerate(image_files):
                    try:
                        # Extract image data
                        image_data = docx_zip.read(image_file)
                        pil_image = Image.open(BytesIO(image_data))
                        
                        # Skip very small images
                        if pil_image.width < 50 or pil_image.height < 50:
                            logger.debug(f"Skipping small image {img_index}: {pil_image.size}")
                            continue
                        
                        images.append(pil_image)
                        logger.info(f"Extracted image from {image_file}: {pil_image.size}")
                        
                    except Exception as e:
                        logger.error(f"Error extracting image {image_file}: {e}")
                        continue
            
            # Method 2: Using python-docx library (as backup)
            if not images and DOCX_IMAGE_SUPPORT:
                try:
                    from docx.document import Document as DocxDoc
                    doc = DocxDoc(file_path)
                    
                    # Check for inline shapes (images)
                    for paragraph in doc.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, '_element') and run._element.xpath('.//a:blip'):
                                logger.info("Found inline image via python-docx")
                                # This is more complex to extract, but the ZIP method above should catch it
                    
                except Exception as e:
                    logger.debug(f"python-docx method failed: {e}")
                    
        except Exception as e:
            logger.error(f"Error extracting images from .docx file: {e}")
        
        return images
    
    def _extract_images_from_xlsx(self, file_path):
        """
        Extract embedded images from Excel files using openpyxl
        
        Args:
            file_path: Path to the .xlsx file
            
        Returns:
            List of PIL Images extracted from the workbook
        """
        images = []
        
        try:
            # Try to import openpyxl
            try:
                from openpyxl import load_workbook
                from openpyxl.drawing.image import Image as OpenpyxlImage
            except ImportError:
                logger.warning("openpyxl not available for Excel image extraction. Install with: pip install openpyxl")
                return images
            
            # Load the workbook
            workbook = load_workbook(file_path)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Check if sheet has any images
                if hasattr(sheet, '_images') and sheet._images:
                    logger.info(f"Found {len(sheet._images)} images in sheet '{sheet_name}'")
                    
                    for img_index, image in enumerate(sheet._images):
                        try:
                            # Get image data
                            if hasattr(image, 'ref') and hasattr(image, '_data'):
                                # Extract image data
                                image_data = image._data()
                                pil_image = Image.open(BytesIO(image_data))
                                
                                # Skip very small images
                                if pil_image.width < 50 or pil_image.height < 50:
                                    logger.debug(f"Skipping small image {img_index}: {pil_image.size}")
                                    continue
                                
                                images.append(pil_image)
                                logger.info(f"Extracted image {img_index} from sheet '{sheet_name}': {pil_image.size}")
                                
                        except Exception as e:
                            logger.error(f"Error extracting image {img_index} from sheet '{sheet_name}': {e}")
                            continue
            
            workbook.close()
            
        except Exception as e:
            logger.error(f"Error extracting images from Excel file: {e}")
        
        return images
    
    def _convert_csv_to_xlsx(self, csv_file_path):
        """
        Convert CSV file to XLSX format
        
        Args:
            csv_file_path: Path to the CSV file
            
        Returns:
            Path to the converted XLSX file
        """
        try:
            import pandas as pd
            
            # Read CSV file
            df = pd.read_csv(csv_file_path)
            
            # Create XLSX file path
            xlsx_file_path = csv_file_path.replace('.csv', '_converted.xlsx')
            
            # Write to Excel
            df.to_excel(xlsx_file_path, index=False, engine='openpyxl')
            
            logger.info(f"Converted CSV to XLSX: {csv_file_path} -> {xlsx_file_path}")
            return xlsx_file_path
            
        except Exception as e:
            logger.error(f"Error converting CSV to XLSX: {e}")
            return None
    
    def _analyze_extracted_images(self, images, source_info):
        """
        Analyze extracted images and create document chunks
        
        Args:
            images: List of PIL Images
            source_info: Dict with source, page/slide info
            
        Returns:
            List of Document chunks with image summaries
        """
        chunks = []
        
        # Create extracted_images folder if it doesn't exist
        extracted_images_folder = "extracted_images"
        os.makedirs(extracted_images_folder, exist_ok=True)
        
        for i, image in enumerate(images):
            try:
                # Save extracted image to folder for reference
                source_name = os.path.basename(source_info["source"]).replace('.', '_')
                page_info = source_info.get("page", 1)
                image_filename = f"{source_name}_page_{page_info}_image_{i}.png"
                image_path = os.path.join(extracted_images_folder, image_filename)
                
                # Save the extracted embedded image
                image.save(image_path)
                logger.info(f"Saved extracted embedded image: {image_path}")
                
                # Analyze image with Azure OpenAI Vision
                image_summary = self.image_analyzer.analyze_image_with_llm(image)
                
                if image_summary and image_summary.strip():
                    # Convert image to base64 for storage
                    import io
                    import base64
                    
                    buffer = io.BytesIO()
                    image.save(buffer, format='PNG')
                    image_base64 = base64.b64encode(buffer.getvalue()).decode()
                    
                    doc_chunk = Document(
                        page_content=f"[EXTRACTED IMAGE SUMMARY] {image_summary}",
                        metadata={
                            "source": source_info["source"],
                            "chunk_id": f"{source_info['source']}_extracted_image_{i}",
                            "chunk_index": source_info.get("chunk_index", i),
                            "page": source_info.get("page", 1),
                            "content_type": "image_summary",
                            "image_type": "extracted_embedded_image",
                            "image_size": f"{image.width}x{image.height}",
                            "image_base64": image_base64,  # Store the actual image
                            "saved_image_path": image_path  # Reference to saved file
                        }
                    )
                    chunks.append(doc_chunk)
                    logger.info(f"Created summary for extracted image {i}: {image.size}")
                    
            except Exception as e:
                logger.error(f"Error analyzing extracted image {i}: {e}")
                continue
                
        return chunks


class RAGPipeline:
    """Complete RAG Pipeline with multimodal document processing"""
    
    def __init__(self, cohere_api_key=None):
        """
        Initialize the RAG pipeline
        
        Args:
            cohere_api_key: API key for Cohere reranking
        """
        self.cohere_api_key = cohere_api_key or os.getenv("COHERE_API_KEY")
        
        # Initialize components
        self.image_analyzer = ImageAnalyzer()
        self.document_processor = DocumentProcessor(self.image_analyzer)
        
        # Initialize embeddings
        self.embeddings = AzureOpenAIEmbeddings(
            azure_deployment=os.getenv("AZURE_EMBEDDING_DEPLOYMENT"),
            api_version=os.getenv("AZURE_EMBEDDING_API_VERSION"),
            api_key=os.getenv("AZURE_EMBEDDING_API_KEY"),
            azure_endpoint=os.getenv("AZURE_EMBEDDING_ENDPOINT")
        )
        
        # Initialize vector store (will be set after adding documents)
        self.vectorstore = None
        self.retriever = None
        
        # Initialize LLM for answer generation
        self.llm = AzureChatOpenAI(
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT"),
            api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
            temperature=0
        )
        
        logger.info("RAG Pipeline initialized successfully")
    
    def add_documents(self, file_paths):
        """
        Process and add documents to the vector store
        
        Args:
            file_paths: List of file paths to process
        """
        logger.info(f"Processing {len(file_paths)} documents...")
        
        # Process documents
        documents, visual_metadata = self.document_processor.process_documents(file_paths)
        
        if not documents:
            raise ValueError("No documents were processed successfully")
        
        logger.info(f"Processed {len(documents)} document chunks")
        
        # Create vector store
        self.vectorstore = Chroma.from_documents(
            documents=documents,
            embedding=self.embeddings,
            persist_directory="./chroma_db"
        )
        
        # Create retriever (will be updated dynamically based on fetch_k)
        self.retriever = self.vectorstore.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 10}  # Default, will be overridden in query method
        )
        
        logger.info("Documents indexed successfully")
        
        return len(documents)
    
    def query(self, question, fetch_k=50, rerank_k=10, use_reranker=True, include_visual_citation=True):
        """
        Query the RAG pipeline with proper retrieval and reranking flow
        
        Flow:
        1. Initial Retrieval: Fetch 'fetch_k' documents from vector store using similarity search
        2. Optional Reranking: If use_reranker=True, send all 'fetch_k' docs to Cohere reranker
        3. Final Selection: Reranker returns top 'rerank_k' documents, or take top 'rerank_k' from fetch_k if no reranking
        4. Context Generation: Use final documents to generate LLM response
        
        Args:
            question: User question
            fetch_k: Number of documents to fetch initially from vector store (default: 50)
            rerank_k: Number of documents after reranking to send to LLM (default: 10)
            use_reranker: Whether to use Cohere reranking (default: True)
            include_visual_citation: Whether to include visual citations (default: True)
            
        Returns:
            Dict with answer, sources, and metadata
        """
        if not self.vectorstore:
            raise ValueError("No documents have been added to the pipeline")
        
        logger.info("="*60)
        logger.info("QUERY PROCESSING STARTED")
        logger.info(f"Query: {question}")
        logger.info(f"Parameters: fetch_k={fetch_k}, rerank_k={rerank_k}, use_reranker={use_reranker}")
        logger.info("="*60)
        
        # Step 1: Initial retrieval using fetch_k
        logger.info("Step 1: Initial Retrieval")
        # Update retriever to use the specified fetch_k
        self.retriever.search_kwargs["k"] = fetch_k
        retrieved_docs = self.retriever.get_relevant_documents(question)
        logger.info(f"Retrieved {len(retrieved_docs)} documents from vector store (fetch_k={fetch_k})")
        
        # Step 2: Reranking (only if use_reranker=True)
        if use_reranker and self.cohere_api_key:
            try:
                logger.info("Step 2: Cohere Reranking")
                logger.info(f"Sending {len(retrieved_docs)} documents to Cohere reranker...")
                
                co = cohere.Client(self.cohere_api_key)
                
                # Prepare documents for reranking
                docs_for_reranking = [doc.page_content for doc in retrieved_docs]
                
                # Rerank documents
                rerank_response = co.rerank(
                    model="rerank-english-v3.0",
                    query=question,
                    documents=docs_for_reranking,
                    top_n=rerank_k
                )
                
                # Get reranked documents in order
                final_docs = []
                for result in rerank_response.results:
                    original_doc = retrieved_docs[result.index]
                    # Add relevance score to metadata for reference
                    original_doc.metadata["rerank_score"] = result.relevance_score
                    final_docs.append(original_doc)
                
                logger.info(f"Reranked to top {len(final_docs)} documents (rerank_k={rerank_k})")
                
                # Log rerank scores for debugging
                logger.info("Rerank Scores:")
                for i, doc in enumerate(final_docs):
                    score = doc.metadata.get("rerank_score", 0.0)
                    source = os.path.basename(doc.metadata.get("source", "Unknown"))
                    page = doc.metadata.get("page", "?")
                    logger.info(f"      Rank {i+1}: Score {score:.4f} | {source} (Page {page}) | {doc.page_content[:80]}...")
                
            except Exception as e:
                logger.warning(f"Reranking failed: {e}")
                logger.info("Falling back to similarity-based selection")
                final_docs = retrieved_docs[:rerank_k]
        else:
            # Step 3: No reranking - just take top rerank_k documents from fetch_k
            logger.info("Step 2: No Reranking (Skipped)")
            final_docs = retrieved_docs[:rerank_k]
            logger.info(f"Using top {len(final_docs)} documents from initial retrieval")
        
        # Prepare context
        context = "\n\n".join([doc.page_content for doc in final_docs])
        logger.info("Step 3: Context Preparation")
        logger.info(f"Context length: {len(context)} characters")
        logger.info(f"Context preview: {context[:200]}..." if len(context) > 200 else f"Full context: {context}")
        
        # Generate answer with improved prompt for visual content
        # Check if we have image summaries in the context
        has_image_content = any("[EXTRACTED IMAGE SUMMARY]" in doc.page_content for doc in final_docs)
        
        if has_image_content:
            prompt = f"""Based on the following context, answer the question. The context includes both text content and image descriptions from documents. Image descriptions are marked with "[EXTRACTED IMAGE SUMMARY]" and contain visual information that may be relevant to the question.

When answering:
- Use information from both text content and image descriptions
- If the question asks about visual elements (restaurants, locations, people, etc.), pay special attention to image descriptions
- Provide specific details from the images when relevant
- If the answer cannot be found in the context, say "I don't have enough information to answer this question."

Context:
{context}

Question: {question}

Answer:"""
        else:
            prompt = f"""Based on the following context, answer the question. If the answer cannot be found in the context, say "I don't have enough information to answer this question."

Context:
{context}

Question: {question}

Answer:"""
        
        try:
            response = self.llm.invoke(prompt)
            answer = response.content
            logger.info("Step 3: LLM Answer Generation")
            logger.info("Generated Answer:")
            logger.info(f"   {answer}")
        except Exception as e:
            logger.error(f"Error generating answer: {e}")
            answer = "I encountered an error while generating the answer."
            logger.info("Step 3: LLM Answer Generation")
            logger.info("Generated Answer:")
            logger.info(f"   {answer}")
        
        # Prepare sources (only show top 1-2 most relevant sources that actually contributed to the answer)
        sources = []
        
        if final_docs:
            # Only take the top 1-2 most relevant documents (these are already ranked by relevance)
            top_docs = final_docs[:2]  # Maximum 2 sources
            seen_sources = set()
            
            for doc in top_docs:
                source_file = doc.metadata.get("source", "Unknown")
                source_page = doc.metadata.get("page", "Unknown")
                content_type = doc.metadata.get("content_type", "text")
                
                # Create a unique identifier for this source+page combination
                source_key = f"{source_file}::{source_page}"
                
                if source_key not in seen_sources:
                    seen_sources.add(source_key)
                    
                    # Clean up the source path to show just the filename
                    source_filename = os.path.basename(source_file) if source_file != "Unknown" else "Unknown"
                    
                    sources.append({
                        "source": source_filename,
                        "page": source_page,
                        "content_type": content_type,
                        "full_path": source_file,  # Keep full path for reference
                        "relevance_rank": len(sources) + 1  # Track relevance ranking
                    })
            
            # If we only have 1 unique source but multiple pages, show max 2 pages
            if len(sources) == 1 and len(final_docs) > 1:
                # Check if the second document is from the same file but different page
                second_doc = final_docs[1]
                second_source_file = second_doc.metadata.get("source", "Unknown")
                second_source_page = second_doc.metadata.get("page", "Unknown")
                second_content_type = second_doc.metadata.get("content_type", "text")
                
                second_source_key = f"{second_source_file}::{second_source_page}"
                
                if (second_source_key not in seen_sources and 
                    second_source_file == sources[0]["full_path"]):  # Same file, different page
                    
                    second_source_filename = os.path.basename(second_source_file) if second_source_file != "Unknown" else "Unknown"
                    
                    sources.append({
                        "source": second_source_filename,
                        "page": second_source_page,
                        "content_type": second_content_type,
                        "full_path": second_source_file,
                        "relevance_rank": 2
                    })
        
        # Prepare visual citations if requested
        visual_citations = []
        if include_visual_citation:
            # Check if query might be related to visual content
            visual_keywords = ["image", "picture", "photo", "chart", "graph", "diagram", "logo", "restaurant", "location", "building", "person", "team", "leadership", "office", "product", "design"]
            is_visual_query = any(keyword.lower() in question.lower() for keyword in visual_keywords)
            
            if is_visual_query:
                # Look for image documents in the results
                for doc in final_docs:
                    if (doc.metadata.get("content_type") == "image_summary" and 
                        doc.metadata.get("image_base64")):
                        
                        visual_citations.append({
                            "source": doc.metadata.get("source", "Unknown"),
                            "page": doc.metadata.get("page", "Unknown"),
                            "chunk_text": doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content,
                            "image_base64": doc.metadata.get("image_base64"),
                            "image_size": doc.metadata.get("image_size", "Unknown"),
                            "saved_image_path": doc.metadata.get("saved_image_path"),
                            "relevance_score": getattr(doc, 'relevance_score', 0.0)
                        })
                
                if visual_citations:
                    logger.info(f"Found {len(visual_citations)} visual citations for the query")
        
        result = {
            "answer": answer,
            "sources": sources,
            "retrieved_chunks": len(final_docs),
            "visual_citations": visual_citations
        }
        
        logger.info("="*60)
        logger.info("QUERY PROCESSING COMPLETED")
        logger.info(f"Generated answer with {len(final_docs)} context chunks")
        logger.info(f"Primary sources: {len(sources)}")
        logger.info(f"Visual citations: {len(visual_citations)}")
        logger.info("="*60)
        
        return result

if __name__ == "__main__":
    def display_results(result):
        """Display query results including visual citations"""
        print(f"\nAnswer: {result['answer']}")
        print(f"\nRetrieved {result['retrieved_chunks']} relevant chunks")
        
        if result['sources']:
            print(f"\nPrimary Sources ({len(result['sources'])} most relevant):")
            for i, source in enumerate(result['sources'], 1):
                content_type_display = f" [{source['content_type']}]" if source['content_type'] != 'text' else ""
                page_display = f"Page {source['page']}" if source['page'] != "Unknown" else "Unknown page"
                relevance_indicator = " Most Relevant" if i == 1 else "Secondary"
                print(f"   {relevance_indicator}: {source['source']} ({page_display}){content_type_display}")
        
        if result['visual_citations']:
            print(f"\nVisual Citations Found: {len(result['visual_citations'])}")
            print("="*60)
            
            for i, citation in enumerate(result['visual_citations'], 1):
                print(f"\nImage {i}:")
                print(f"   Source: {citation['source']}")
                print(f"   Page: {citation['page']}")
                print(f"   Size: {citation['image_size']}")
                print(f"   Description: {citation['chunk_text']}")
                
                # Show saved image path if available
                if 'saved_image_path' in citation:
                    print(f"   Extracted Image: {citation['saved_image_path']}")
                
                # Save image to file for viewing
                if citation.get('image_base64'):
                    try:
                        import base64
                        from PIL import Image
                        import io
                        
                        # Decode base64 image
                        image_data = base64.b64decode(citation['image_base64'])
                        image = Image.open(io.BytesIO(image_data))
                        
                        # Save image with a descriptive filename
                        source_clean = citation['source'].replace('/', '_').replace('\\', '_')
                        filename = f"extracted_image_{i}_{source_clean}_page_{citation['page']}.png"
                        image.save(filename)
                        
                    except Exception as e:
                        print(f"   Error saving image: {e}")
                
                print("-" * 40)
        
        print("\n" + "="*50)

    def main():
        print("RAG Pipeline")
        print("=" * 60)
        print(f"Logs are being saved to: {log_filename}")
        print("=" * 60)

        # --- Configuration ---
        COHERE_API_KEY = os.getenv("COHERE_API_KEY")
        
        # Get all files from the files folder
        files_folder = "files"
        supported_extensions = ['.pdf', '.docx', '.txt', '.pptx', '.xlsx']
        
        if not os.path.exists(files_folder):
            print(f"ERROR: Files folder '{files_folder}' not found!")
            return
        
        file_paths = []
        for filename in os.listdir(files_folder):
            file_path = os.path.join(files_folder, filename)
            if os.path.isfile(file_path):
                # Skip temporary files
                if filename.startswith('~$') or filename.startswith('.~'):
                    continue
                    
                file_ext = os.path.splitext(filename)[1].lower()
                if file_ext in supported_extensions:
                    file_paths.append(file_path)
        
        if not file_paths:
            print(f"ERROR: No supported files found in '{files_folder}' folder!")
            print(f"Supported formats: {', '.join(supported_extensions)}")
            return
        
        # user_query = "Show me information about restaurants or any images related to dining locations"
        # user_query = "What are the formats they're looking for?"
        # user_query = "What's the check-in time and seat number?"
        # user_query = "What is the source information of the ACEN Corp?"
        user_query = "What is q,k and v in the context?"

        print("\nAPI Configuration")
        print("SUCCESS: Cohere API key loaded from existing configuration")

        # Initialize RAG pipeline
        print("\nInitializing RAG Pipeline...")
        try:
            rag = RAGPipeline(
                cohere_api_key=COHERE_API_KEY
            )
            print("SUCCESS: Pipeline initialized successfully!")
        except Exception as e:
            print(f"ERROR: Failed to initialize pipeline: {e}")
            return

        print("\nDocument Processing")
        print(f"Found {len(file_paths)} supported file(s) in '{files_folder}' folder:")
        for i, file_path in enumerate(file_paths, 1):
            print(f"  {i}. {os.path.basename(file_path)}")

        # Process documents
        print(f"\nProcessing {len(file_paths)} documents...")
        try:
            rag.add_documents(file_paths)
            print("SUCCESS: Documents processed and indexed successfully!")
        except Exception as e:
            print(f"ERROR: Failed to process documents: {e}")
            return

        # Query
        print("\n" + "="*60)
        print("PROCESSING QUERY")
        print("="*60)
        print(f"User query: {user_query}")

        # Demonstration of corrected retrieval flow
        print(f"\nRetrieval Flow:")
        print(f"1. Initial retrieval: fetch_k=50 documents from vector store")
        print(f"2. Reranking: use_reranker=True - send all 50 docs to Cohere reranker")
        print(f"3. Final selection: rerank_k=10 - reranker returns top 10 documents")
        print(f"4. Context generation: Use top 10 documents for LLM response")

        try:
            result = rag.query(
                question=user_query, 
                fetch_k=50,      # Retrieve 50 documents initially from vector store
                rerank_k=10,     # Rerank and select top 10 documents for LLM context
                use_reranker=True,  # Use Cohere reranking after initial retrieval
                include_visual_citation=True
            )
            print("\n" + "="*60)
            print("RESULTS")
            print("="*60)
            display_results(result)
        except Exception as e:
            print(f"ERROR: Error processing query: {e}")

        print("\nThanks for using the RAG Pipeline!")

    # Run the main function
    main()
